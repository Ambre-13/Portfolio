"""
Convertisseur Markdown → PowerPoint ultra complet
- Format 16:9 widescreen
- Thème Azure (bleu professionnel)
- Tout le texte affiché sans rien tronquer
- Titre + sous-titre si '#' et '##' présents
- Listes numérotées, bullets imbriqués, texte libre
- Taille automatique du texte si débordement
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

# ── Chemins ──────────────────────────────────────────────────────────────────
MD_PATH   = Path("Diaporama_Azure_Complet.md")
PPTX_PATH = Path("Diaporama_Azure_Complet.pptx")

# ── Couleurs thème Azure ──────────────────────────────────────────────────────
BLEU_AZURE      = RGBColor(0x00, 0x78, 0xD4)
BLEU_FONCE      = RGBColor(0x00, 0x35, 0x6B)
BLANC           = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_TEXTE      = RGBColor(0x21, 0x21, 0x21)
ORANGE_ACCENT   = RGBColor(0xE9, 0x71, 0x32)

# ── Dimensions 16:9 ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.45)


# ══════════════════════════════════════════════════════════════════════════════
# Parsing du Markdown
# ══════════════════════════════════════════════════════════════════════════════

def strip_front_matter(text: str) -> str:
    lines = text.splitlines()
    if lines and lines[0].strip() == "---":
        for i in range(1, len(lines)):
            if lines[i].strip() == "---":
                return "\n".join(lines[i + 1:])
    return text


def split_slides(text: str) -> list[str]:
    parts = re.split(r"(?m)^---\s*$", text)
    return [p.strip() for p in parts if p.strip()]


def clean_md(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*",     r"\1", text)
    text = re.sub(r"`(.*?)`",       r"\1", text)
    text = re.sub(r"\$(.*?)\$",     r"\1", text)
    text = re.sub(r"\\\w+\$?",      "",    text)
    return text.strip()


def is_bold_md(raw: str) -> bool:
    return bool(re.search(r"\*\*(.+?)\*\*", raw))


class Item:
    __slots__ = ("text", "kind", "level", "bold")

    def __init__(self, text: str, kind: str = "text",
                 level: int = 0, bold: bool = False):
        self.text  = text
        self.kind  = kind   # title2 | bullet | numbered | text
        self.level = level
        self.bold  = bold


def parse_slide_content(raw: str) -> tuple[str, str, list[Item]]:
    title1 = ""
    title2 = ""
    items: list[Item] = []

    for line in raw.splitlines():
        # H1 → titre principal
        m = re.match(r"^#\s+(.+)$", line)
        if m:
            if not title1:
                title1 = clean_md(m.group(1))
            else:
                items.append(Item(clean_md(m.group(1)), "title2", bold=True))
            continue

        # H2 → sous-titre
        m = re.match(r"^##\s+(.+)$", line)
        if m:
            if not title2 and not items:
                title2 = clean_md(m.group(1))
            else:
                items.append(Item(clean_md(m.group(1)), "title2", bold=True))
            continue

        # H3-H6 → corps gras
        m = re.match(r"^#{3,6}\s+(.+)$", line)
        if m:
            items.append(Item(clean_md(m.group(1)), "text", bold=True))
            continue

        # Bullet
        m = re.match(r"^(\s*)[-*]\s+(.+)$", line)
        if m:
            indent = len(m.group(1).replace("\t", "    "))
            lvl = min(indent // 2, 4)
            items.append(Item(clean_md(m.group(2)), "bullet", level=lvl,
                              bold=is_bold_md(m.group(2))))
            continue

        # Numéroté
        m = re.match(r"^(\s*)\d+[.)]\s+(.+)$", line)
        if m:
            indent = len(m.group(1).replace("\t", "    "))
            lvl = min(indent // 2, 4)
            items.append(Item(clean_md(m.group(2)), "numbered", level=lvl,
                              bold=is_bold_md(m.group(2))))
            continue

        stripped = line.strip()
        if not stripped:
            continue

        items.append(Item(clean_md(stripped), "text",
                          bold=is_bold_md(stripped)))

    return title1, title2, items


# ══════════════════════════════════════════════════════════════════════════════
# Helpers PPTX
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    return tb


def enable_autofit(tf):
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    # Supprimer balises autofit existantes
    for tag in (qn("a:normAutofit"), qn("a:spAutoFit"), qn("a:noAutofit")):
        for child in list(bodyPr):
            if child.tag == tag:
                bodyPr.remove(child)
    bodyPr.set("wrap", "square")
    norm = etree.SubElement(bodyPr, qn("a:normAutofit"))
    norm.set("fontScale", "100000")
    norm.set("lnSpcReduction", "20000")


def set_para(tf, text: str, font_sz: float, color: RGBColor,
             bold: bool = False, first: bool = False,
             align=PP_ALIGN.LEFT, space_before: int = 0):
    """Ajoute un paragraphe dans le text_frame."""
    if first:
        para = tf.paragraphs[0]
        para.clear()
    else:
        para = tf.add_paragraph()

    para.alignment = align
    if space_before:
        para.space_before = Pt(space_before)

    run = para.add_run()
    run.text = text
    run.font.size  = Pt(font_sz)
    run.font.color.rgb = color
    run.font.bold  = bold
    return para


def compute_font_size(n: int) -> float:
    if n <= 7:  return 18.0
    if n <= 12: return 16.0
    if n <= 18: return 14.0
    if n <= 26: return 12.5
    if n <= 36: return 11.0
    return 10.0


# ══════════════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════════════

def build_cover_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond dégradé simulé (deux rect)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLEU_FONCE)
    add_rect(slide, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), ORANGE_ACCENT)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, ORANGE_ACCENT)

    # Titre
    tb = add_textbox(slide, MARGIN + Inches(0.2), Inches(1.6),
                     SLIDE_W - MARGIN * 2 - Inches(0.2), Inches(2.8))
    tf = tb.text_frame
    tf.word_wrap = True

    set_para(tf, title, 38, BLANC, bold=True, first=True, align=PP_ALIGN.LEFT)

    # Sous-titre
    if subtitle:
        set_para(tf, subtitle, 22, RGBColor(0xC8, 0xE0, 0xFF),
                 first=False, align=PP_ALIGN.LEFT, space_before=8)

    # Label bas
    tb2 = add_textbox(slide, MARGIN, SLIDE_H - Inches(0.6),
                      SLIDE_W - MARGIN * 2, Inches(0.45))
    tf2 = tb2.text_frame
    set_para(tf2, "Microsoft Azure — Formation Bloc 2", 13,
             RGBColor(0xA0, 0xC4, 0xEE), first=True, align=PP_ALIGN.LEFT)


def build_content_slide(prs: Presentation, title1: str, title2: str,
                        items: list[Item]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLANC)

    # Barre titre bleue
    header_h = Inches(1.15)
    add_rect(slide, 0, 0, SLIDE_W, header_h, BLEU_AZURE)

    # Trait orange
    add_rect(slide, 0, header_h, SLIDE_W, Inches(0.05), ORANGE_ACCENT)

    # Texte titre
    tb_t = add_textbox(slide, MARGIN, Inches(0.08),
                       SLIDE_W - MARGIN * 2, header_h - Inches(0.08))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    t1_display = title1 or "Contenu"
    set_para(tf_t, t1_display, 24, BLANC, bold=True, first=True, align=PP_ALIGN.LEFT)

    if title2:
        set_para(tf_t, title2, 15, RGBColor(0xC8, 0xE0, 0xFF),
                 first=False, align=PP_ALIGN.LEFT)

    # Corps
    if not items:
        return

    body_top = header_h + Inches(0.2)
    body_h   = SLIDE_H - body_top - Inches(0.18)

    tb_b = add_textbox(slide, MARGIN, body_top,
                       SLIDE_W - MARGIN * 2, body_h)
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    enable_autofit(tf_b)

    fs = compute_font_size(len(items))

    BULLET_PREFIXES = {0: "• ", 1: "  ◦ ", 2: "    ▪ ", 3: "      – ", 4: "        · "}
    NUM_INDENT      = {0: "", 1: "  ", 2: "    ", 3: "      ", 4: "        "}

    first_p = True
    for item in items:
        if item.kind == "bullet":
            prefix = BULLET_PREFIXES.get(item.level, "• ")
            text_display = prefix + item.text
            color  = GRIS_TEXTE
            size   = fs
        elif item.kind == "numbered":
            prefix = NUM_INDENT.get(item.level, "")
            text_display = prefix + item.text
            color  = GRIS_TEXTE
            size   = fs
        elif item.kind == "title2":
            text_display = item.text
            color  = BLEU_AZURE
            size   = max(fs, 14.0)
        elif item.bold:
            text_display = item.text
            color  = BLEU_FONCE
            size   = fs
        else:
            text_display = item.text
            color  = GRIS_TEXTE
            size   = fs

        if first_p:
            para = tf_b.paragraphs[0]
            para.clear()
            first_p = False
        else:
            para = tf_b.add_paragraph()

        para.alignment = PP_ALIGN.LEFT

        run = para.add_run()
        run.text = text_display
        run.font.size      = Pt(size)
        run.font.color.rgb = color
        run.font.bold      = item.bold if item.kind not in ("bullet", "numbered") else False


# ══════════════════════════════════════════════════════════════════════════════
# Entrée principale
# ══════════════════════════════════════════════════════════════════════════════

def main():
    if not MD_PATH.exists():
        raise FileNotFoundError(f"Introuvable : {MD_PATH}")

    md = MD_PATH.read_text(encoding="utf-8")
    md = strip_front_matter(md)
    slides_raw = split_slides(md)
    print(f"Slides détectées dans le Markdown : {len(slides_raw)}")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, raw in enumerate(slides_raw):
        t1, t2, items = parse_slide_content(raw)

        if idx == 0:
            cover_title = t1 if t1 else "Bloc 2 – Azure"
            cover_sub   = t2 if t2 else ""
            if items and not cover_sub:
                cover_sub = "  |  ".join(
                    it.text for it in items[:4] if it.text)
            build_cover_slide(prs, cover_title, cover_sub)
        else:
            build_content_slide(prs, t1, t2, items)

    prs.save(PPTX_PATH)

    total = len(prs.slides)
    size_kb = PPTX_PATH.stat().st_size // 1024
    print(f"✅  PowerPoint généré : {PPTX_PATH.resolve()}")
    print(f"   Slides : {total}   |   Taille : {size_kb} Ko")


if __name__ == "__main__":
    main()
